"""SDN3仮想化基盤開発 提案書 PPTX生成スクリプト

story.md + HTMLデザインモックに基づき、全34スライドのPPTXを生成する。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Optional

# ============================================================
# テーマ定義
# ============================================================

# 配色（HTMLデザインモックから抽出）
C = {
    "navy_900": RGBColor(0x06, 0x12, 0x25),
    "navy_800": RGBColor(0x0B, 0x1D, 0x3A),
    "navy_700": RGBColor(0x10, 0x2A, 0x4F),
    "navy_600": RGBColor(0x16, 0x3D, 0x6E),
    "blue_500": RGBColor(0x00, 0x67, 0xC5),
    "blue_400": RGBColor(0x2B, 0x8A, 0xDB),
    "blue_300": RGBColor(0x5A, 0xAD, 0xEE),
    "teal_500": RGBColor(0x00, 0xB4, 0xD8),
    "teal_400": RGBColor(0x22, 0xD3, 0xEE),
    "green_500": RGBColor(0x22, 0xC5, 0x5E),
    "amber_500": RGBColor(0xF5, 0x9E, 0x0B),
    "red_500": RGBColor(0xEF, 0x44, 0x44),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "bg_light": RGBColor(0xF8, 0xFA, 0xFC),
    "bg_alt": RGBColor(0xF1, 0xF5, 0xF9),
    "border": RGBColor(0xE2, 0xE8, 0xF0),
    "text_primary": RGBColor(0x0F, 0x17, 0x2A),
    "text_secondary": RGBColor(0x47, 0x55, 0x69),
    "text_muted": RGBColor(0x94, 0xA3, 0xB8),
}

# フォント
FONT_HEADING = "メイリオ"
FONT_BODY = "メイリオ"

# スライドサイズ (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# レイアウト定数
MARGIN_LEFT = Inches(0.7)
MARGIN_RIGHT = Inches(0.7)
MARGIN_TOP = Inches(0.5)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
TOPBAR_HEIGHT = Inches(0.06)
FOOTER_Y = Inches(7.0)

# 出力先
BASE_DIR = os.path.dirname(__file__)
ASSETS_DIR = os.path.join(BASE_DIR, "assets")
OUTPUT_DIR = os.path.join(BASE_DIR, "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)


# ============================================================
# ヘルパー関数
# ============================================================

def _set_font(
    run,
    name: str = FONT_BODY,
    size: Optional[Pt] = None,
    bold: bool = False,
    color: Optional[RGBColor] = None,
) -> None:
    """ランのフォント設定を行う。

    :param run: python-pptx Run オブジェクト
    :param name: フォント名
    :param size: フォントサイズ
    :param bold: 太字
    :param color: フォント色
    """
    run.font.name = name
    if size:
        run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_name: str = FONT_BODY,
    font_size: Pt = Pt(14),
    bold: bool = False,
    color: RGBColor = None,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
    vertical: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    """テキストボックスを追加する。

    :param slide: スライドオブジェクト
    :param left: 左位置
    :param top: 上位置
    :param width: 幅
    :param height: 高さ
    :param text: テキスト
    :param font_name: フォント名
    :param font_size: フォントサイズ
    :param bold: 太字
    :param color: フォント色
    :param alignment: テキスト配置
    :param word_wrap: 折り返し
    :param vertical: 垂直配置
    :return: TextFrameオブジェクト
    """
    if color is None:
        color = C["text_primary"]
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, font_name, font_size, bold, color)
    # vertical anchor
    txBox.text_frame.paragraphs[0].space_before = Pt(0)
    txBox.text_frame.paragraphs[0].space_after = Pt(0)
    return tf


def _add_rect(
    slide,
    left,
    top,
    width,
    height,
    fill_color: Optional[RGBColor] = None,
    line_color: Optional[RGBColor] = None,
    line_width: Pt = Pt(0),
):
    """矩形シェイプを追加する。

    :param slide: スライドオブジェクト
    :param left: 左位置
    :param top: 上位置
    :param width: 幅
    :param height: 高さ
    :param fill_color: 塗り色
    :param line_color: 線色
    :param line_width: 線幅
    :return: Shapeオブジェクト
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_topbar(slide) -> None:
    """スライド上部のグラデーション風バーを追加する。

    :param slide: スライドオブジェクト
    """
    # グラデーションの代わりに2色分割
    half_w = SLIDE_WIDTH // 2
    _add_rect(slide, Inches(0), Inches(0), half_w, TOPBAR_HEIGHT,
              fill_color=C["blue_500"])
    _add_rect(slide, half_w, Inches(0), half_w, TOPBAR_HEIGHT,
              fill_color=C["teal_500"])


def _add_footer(slide, page_num: int) -> None:
    """フッターを追加する。

    :param slide: スライドオブジェクト
    :param page_num: ページ番号
    """
    # 線
    _add_rect(slide, MARGIN_LEFT, FOOTER_Y, CONTENT_WIDTH, Pt(1),
              fill_color=C["border"])
    # 左テキスト
    _add_textbox(slide, MARGIN_LEFT, FOOTER_Y + Inches(0.05),
                 Inches(5), Inches(0.3),
                 "SDN3仮想化基盤開発 ご提案書",
                 font_size=Pt(8), color=C["text_muted"])
    # 右ページ番号
    _add_textbox(slide, SLIDE_WIDTH - MARGIN_RIGHT - Inches(0.5),
                 FOOTER_Y + Inches(0.05),
                 Inches(0.5), Inches(0.3),
                 str(page_num),
                 font_size=Pt(8), color=C["text_muted"],
                 alignment=PP_ALIGN.RIGHT)


def _add_page_badge(slide, num: int) -> None:
    """ページ番号バッジを追加する。

    :param slide: スライドオブジェクト
    :param num: ページ番号
    """
    badge = _add_rect(slide, MARGIN_LEFT, Inches(0.35), Inches(0.5), Inches(0.3),
                      fill_color=RGBColor(0xE8, 0xF0, 0xFE))
    badge.line.fill.background()
    _add_textbox(slide, MARGIN_LEFT + Inches(0.02), Inches(0.37),
                 Inches(0.46), Inches(0.26),
                 f"{num:02d}",
                 font_name=FONT_HEADING, font_size=Pt(10),
                 bold=True, color=C["blue_500"],
                 alignment=PP_ALIGN.CENTER)


def _add_slide_title(slide, page_num: int, title: str) -> None:
    """スライドタイトルを追加する（バッジ＋タイトル＋バー＋フッター）。

    :param slide: スライドオブジェクト
    :param page_num: ページ番号
    :param title: タイトル
    """
    _add_topbar(slide)
    _add_page_badge(slide, page_num)
    _add_textbox(slide, MARGIN_LEFT + Inches(0.6), Inches(0.3),
                 CONTENT_WIDTH - Inches(0.6), Inches(0.45),
                 title,
                 font_name=FONT_HEADING, font_size=Pt(22),
                 bold=True, color=C["navy_800"])
    _add_footer(slide, page_num)


def _add_message_bar(slide, message: str, top: float = None) -> float:
    """メッセージ帯を追加する。

    :param slide: スライドオブジェクト
    :param message: メッセージテキスト
    :param top: 上位置（省略時はデフォルト）
    :return: メッセージ帯下端のY座標
    """
    if top is None:
        top = Inches(0.85)
    bar_h = Inches(0.4)
    # 背景
    _add_rect(slide, MARGIN_LEFT, top, CONTENT_WIDTH, bar_h,
              fill_color=RGBColor(0xEE, 0xF2, 0xFA))
    # 左ボーダー
    _add_rect(slide, MARGIN_LEFT, top, Inches(0.04), bar_h,
              fill_color=C["blue_500"])
    # テキスト
    _add_textbox(slide, MARGIN_LEFT + Inches(0.15), top + Inches(0.02),
                 CONTENT_WIDTH - Inches(0.2), bar_h - Inches(0.04),
                 message,
                 font_size=Pt(11), bold=True, color=C["navy_700"])
    return top + bar_h + Inches(0.1)


def _add_bullet_card(
    slide,
    left,
    top,
    width,
    height,
    title: str,
    body: str,
    accent_color: RGBColor = None,
    icon_text: str = "",
) -> None:
    """カード型の箇条書き項目を追加する。

    :param slide: スライドオブジェクト
    :param left: 左位置
    :param top: 上位置
    :param width: 幅
    :param height: 高さ
    :param title: 見出し
    :param body: 本文
    :param accent_color: アクセントカラー
    :param icon_text: アイコンテキスト
    """
    if accent_color is None:
        accent_color = C["blue_500"]
    # 背景
    _add_rect(slide, left, top, width, height,
              fill_color=C["bg_light"], line_color=C["border"], line_width=Pt(0.5))
    # 左アクセントライン
    _add_rect(slide, left, top, Inches(0.04), height,
              fill_color=accent_color)
    # アイコン
    if icon_text:
        icon = _add_rect(slide, left + Inches(0.15), top + Inches(0.1),
                         Inches(0.35), Inches(0.35),
                         fill_color=accent_color)
        icon.line.fill.background()
        # 角丸
        _add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                     Inches(0.35), Inches(0.35),
                     icon_text,
                     font_name=FONT_HEADING, font_size=Pt(14),
                     bold=True, color=C["white"],
                     alignment=PP_ALIGN.CENTER)
        text_left = left + Inches(0.6)
        text_width = width - Inches(0.75)
    else:
        text_left = left + Inches(0.15)
        text_width = width - Inches(0.3)

    # タイトル
    _add_textbox(slide, text_left, top + Inches(0.08),
                 text_width, Inches(0.25),
                 title,
                 font_name=FONT_HEADING, font_size=Pt(12),
                 bold=True, color=C["navy_800"])
    # 本文
    _add_textbox(slide, text_left, top + Inches(0.32),
                 text_width, height - Inches(0.4),
                 body,
                 font_size=Pt(10), color=C["text_secondary"])


def _add_diagram_placeholder(slide, left, top, width, height, filename: str) -> None:
    """図のプレースホルダーを追加する（PNGが無い場合）。

    :param slide: スライドオブジェクト
    :param left: 左位置
    :param top: 上位置
    :param width: 幅
    :param height: 高さ
    :param filename: ファイル名
    """
    # 背景
    _add_rect(slide, left, top, width, height,
              fill_color=C["bg_light"], line_color=C["border"], line_width=Pt(1))
    # ファイル名ラベル
    _add_textbox(slide, left, top + height // 2 - Inches(0.2),
                 width, Inches(0.4),
                 f"[ {filename} ]",
                 font_size=Pt(11), color=C["text_muted"],
                 alignment=PP_ALIGN.CENTER)


def _try_add_image(slide, left, top, width, height, drawio_name: str) -> None:
    """PNGがあれば画像、なければプレースホルダーを追加する。

    :param slide: スライドオブジェクト
    :param left: 左位置
    :param top: 上位置
    :param width: 幅
    :param height: 高さ
    :param drawio_name: .drawioファイル名（拡張子付き）
    """
    png_name = drawio_name.replace(".drawio", ".png")
    png_path = os.path.join(ASSETS_DIR, png_name)
    if os.path.exists(png_path):
        slide.shapes.add_picture(png_path, left, top, width, height)
    else:
        _add_diagram_placeholder(slide, left, top, width, height, drawio_name)


def _make_table(
    slide,
    left,
    top,
    width,
    row_height,
    headers: list[str],
    rows: list[list[str]],
    col_widths: Optional[list[float]] = None,
):
    """テーブルを追加する。

    :param slide: スライドオブジェクト
    :param left: 左位置
    :param top: 上位置
    :param width: テーブル幅
    :param row_height: 行高さ
    :param headers: ヘッダー行
    :param rows: データ行
    :param col_widths: 列幅の比率リスト
    :return: Tableオブジェクト
    """
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_h = row_height * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, total_h)
    table = table_shape.table

    # 列幅
    if col_widths:
        total_ratio = sum(col_widths)
        for i, ratio in enumerate(col_widths):
            table.columns[i].width = int(width * ratio / total_ratio)

    # ヘッダー
    for ci, header_text in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = header_text
        _set_font(run, FONT_HEADING, Pt(10), True, C["white"])
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["navy_800"]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # データ行
    for ri, row_data in enumerate(rows):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            _set_font(run, FONT_BODY, Pt(9), False, C["text_primary"])
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # 交互背景
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["bg_light"]

    return table


def _make_gantt_cell_colored(table, row: int, col: int, color: RGBColor) -> None:
    """ガントチャートのセルに色を設定する。

    :param table: Tableオブジェクト
    :param row: 行番号
    :param col: 列番号
    :param color: 背景色
    """
    cell = table.cell(row, col)
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


# ============================================================
# スライド生成関数
# ============================================================

def slide_01_cover(prs: Presentation) -> None:
    """スライド1: 表紙。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ダークネイビー背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C["navy_800"]

    # 上部グラデーションバー
    half = SLIDE_WIDTH // 2
    _add_rect(slide, Inches(0), Inches(0), half, Inches(0.06),
              fill_color=C["blue_500"])
    _add_rect(slide, half, Inches(0), half, Inches(0.06),
              fill_color=C["teal_500"])

    # Confidential
    _add_textbox(slide, Inches(1.0), Inches(1.5), Inches(5), Inches(0.3),
                 "FOR YOUR EYES ONLY B",
                 font_name=FONT_HEADING, font_size=Pt(10),
                 bold=True, color=C["teal_400"])

    # タイトル
    _add_textbox(slide, Inches(1.0), Inches(2.2), Inches(8), Inches(1.0),
                 "法人SDN基盤\nSDN3仮想化基盤開発",
                 font_name=FONT_HEADING, font_size=Pt(36),
                 bold=True, color=C["white"])

    # サブタイトル
    _add_textbox(slide, Inches(1.0), Inches(3.5), Inches(5), Inches(0.5),
                 "ご提案書",
                 font_name=FONT_HEADING, font_size=Pt(20),
                 color=C["blue_300"])

    # メタ情報
    meta_items = [
        ("提出日", "2026年2月"),
        ("提出先", "ソフトバンク株式会社 御中"),
    ]
    for i, (label, value) in enumerate(meta_items):
        x = Inches(1.0) + Inches(2.5) * i
        _add_textbox(slide, x, Inches(4.5), Inches(2.3), Inches(0.2),
                     label,
                     font_size=Pt(8), bold=True, color=C["text_muted"])
        _add_textbox(slide, x, Inches(4.75), Inches(2.3), Inches(0.25),
                     value,
                     font_size=Pt(12), bold=True, color=C["white"])

    # 会社名
    _add_textbox(slide, Inches(9.0), Inches(5.8), Inches(3.5), Inches(0.4),
                 "COMPANY NAME",
                 font_name=FONT_HEADING, font_size=Pt(16),
                 bold=True, color=C["white"],
                 alignment=PP_ALIGN.RIGHT)
    _add_textbox(slide, Inches(9.0), Inches(6.2), Inches(3.5), Inches(0.3),
                 "Technology & Integration",
                 font_size=Pt(10), color=C["text_muted"],
                 alignment=PP_ALIGN.RIGHT)


def slide_02_toc(prs: Presentation) -> None:
    """スライド2: 目次。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 2, "目次")
    content_top = _add_message_bar(slide, "提案書の構成")

    sections = [
        ("01", "エグゼクティブサマリー"),
        ("02", "貴社要求の理解"),
        ("03", "ソリューション全体構成"),
        ("04", "ストレージソリューション提案"),
        ("05", "設計方針"),
        ("06", "非機能要件対応"),
        ("07", "機器・ライセンス情報"),
        ("08", "運用・保守"),
        ("09", "プロジェクト管理"),
        ("10", "まとめ"),
    ]

    for i, (num, title) in enumerate(sections):
        row = i // 2
        col = i % 2
        x = MARGIN_LEFT + Inches(0.1) + col * Inches(5.8)
        y = content_top + row * Inches(0.55)

        # 番号
        _add_textbox(slide, x, y, Inches(0.5), Inches(0.4),
                     num,
                     font_name=FONT_HEADING, font_size=Pt(18),
                     bold=True, color=C["blue_500"])
        # タイトル
        _add_textbox(slide, x + Inches(0.55), y + Inches(0.05),
                     Inches(5), Inches(0.35),
                     title,
                     font_name=FONT_HEADING, font_size=Pt(14),
                     bold=True, color=C["navy_800"])
        # 下線
        _add_rect(slide, x, y + Inches(0.45), Inches(5.5), Pt(0.5),
                  fill_color=C["border"])


def slide_03_executive_summary(prs: Presentation) -> None:
    """スライド3: エグゼクティブサマリー。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 3, "エグゼクティブサマリー")
    y = _add_message_bar(slide, "本提案は「信頼性」「拡張性」「運用効率」の3本柱でSDN3基盤を実現する")

    cards = [
        ("1", "信頼性 — 高可用性ストレージ基盤",
         "NetApp ONTAPベースのストレージ基盤で、HA Pair構成による自動フェイルオーバー、RAID-TECによるディスク冗長を実現。通信事業者の仮想化基盤への豊富な導入実績に裏打ちされた信頼性を提供します。",
         C["blue_500"]),
        ("2", "拡張性 — Pod単位のスケーラブル設計",
         "5年間の収容需要を見据えたサイジングに加え、ディスクシェルフ追加→筐体追加→Pod追加の段階的な拡張パスを設計。サービス無停止での柔軟な増設を可能にします。",
         C["teal_500"]),
        ("3", "運用効率 — 統合保守と自動化",
         "24/365対応の統合保守サービスと、Ansible Automation Platformによる運用自動化で、SB様運用チームの負荷を軽減。月次定例報告とハンズオントレーニングで自立運用を支援します。",
         C["green_500"]),
        ("4", "インテグレーション力 — 豊富な構築実績",
         "通信事業者向け仮想化基盤の豊富な構築実績と、NetApp製品の専門知見を活かした確実なインテグレーションを提供します。",
         C["amber_500"]),
    ]

    card_h = Inches(1.1)
    for i, (icon, title, body, color) in enumerate(cards):
        _add_bullet_card(slide, MARGIN_LEFT, y + i * (card_h + Inches(0.08)),
                         CONTENT_WIDTH, card_h, title, body, color, icon)


def slide_04_requirements(prs: Presentation) -> None:
    """スライド4: 貴社要求の理解。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 4, "貴社要求の理解")
    y = _add_message_bar(slide, "RFPの背景・目的を正確に理解し、各課題に対する解決策を提案する")

    headers = ["背景・課題", "目的", "当社提案の対応方針"]
    rows = [
        ["SDN2基盤の老朽化", "安定したサービスレベル維持", "最新NetApp AFF + ONTAP 9.xで高性能・高信頼を確保"],
        ["収容残量の低下", "広帯域化への対応", "5年需要を見据えたサイジング + Pod単位水平拡張"],
        ["広帯域要望の増加", "耐障害性の強化", "全コンポーネント冗長構成 + HA/DRS自動復旧"],
        ["OnePlatform継続不可", "ログ・統計機能の継続", "SDN3内運用サーバ群で独立したログ/統計基盤を構築"],
    ]

    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.5),
                headers, rows, col_widths=[2, 2, 4])


def slide_05_commercial_overview(prs: Presentation) -> None:
    """スライド5: ソリューション全体構成（商用環境）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 5, "ソリューション全体構成（商用環境）")
    y = _add_message_bar(slide, "東京・大阪の商用環境全体像と開発対象範囲を明示")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.2),
                   "slide05_commercial_overview.drawio")


def slide_06_stg_dev_overview(prs: Presentation) -> None:
    """スライド6: ソリューション全体構成（STG・開発環境）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 6, "ソリューション全体構成（STG・開発環境）")
    y = _add_message_bar(slide, "STG・開発環境は商用に準じた構成で、作業手順の確立と開発検証を支える")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.2),
                   "slide06_stg_dev_overview.drawio")


def slide_07_scope(prs: Presentation) -> None:
    """スライド7: 開発対象範囲。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 7, "開発対象範囲")
    y = _add_message_bar(slide, "当社の調達・構築対象範囲を明確化し、責任分界点を明示")

    headers = ["コンポーネント", "調達", "構築", "商用東京", "商用大阪", "STG", "開発", "備考"]
    rows = [
        ["NFS Storage (リソースPod)", "当社", "当社", "2台", "1台", "2台", "2台", "NetApp AFF"],
        ["VM Storage (管理Pod)", "当社", "当社", "1台", "—", "1台", "1台", "NetApp AFF"],
        ["Data Storage (管理Pod)", "当社", "当社", "1台", "—", "1台", "1台", "NetApp AFF"],
        ["Backup Storage", "当社", "当社", "1台", "—", "1台", "1台", "NetApp FAS"],
        ["NSO (Cisco)", "当社", "当社", "2台", "—", "2台", "2台", "Active/Standby"],
        ["ESXi Server", "SB様", "当社", "18+18+N", "18", "N", "N", "Dell R670"],
        ["NW機器", "SB様", "当社", "多数", "多数", "多数", "多数", "Cisco Nexus"],
        ["iPDU", "当社", "当社", "各ラック", "各ラック", "各ラック", "各ラック", "インテリジェントPDU"],
    ]
    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.42),
                headers, rows, col_widths=[2.5, 1, 1, 1, 1, 0.8, 0.8, 1.8])


def slide_08_storage_overview(prs: Presentation) -> None:
    """スライド8: ストレージ提案概要。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 8, "ストレージ提案概要")
    y = _add_message_bar(slide, "NetApp AFF/FASシリーズで全ストレージ要件を統一的にカバー")

    cards = [
        ("", "リソースPod向け NFS Storage (NetApp AFF)",
         "VM領域に最適化したNFSストレージ。HA Pair構成でコントローラ冗長、RAID-TECでディスク冗長。商用東京2台、大阪1台、STG/開発各2台を配置。",
         C["blue_500"]),
        ("", "管理Pod向け VM/Data Storage (NetApp AFF)",
         "管理VM用NFS領域とファイル共有/ログ保存用データ領域を分離配置。SnapMirrorによるSDN2⇔SDN3ミラーリング構成に対応。",
         C["teal_500"]),
        ("", "バックアップ用 Backup Storage (NetApp FAS)",
         "SnapVaultによる日次バックアップ先。コストパフォーマンスに優れたFASシリーズで大容量を確保。",
         C["green_500"]),
        ("", "統一運用のメリット",
         "ONTAPベースの統一アーキテクチャにより、運用手順の共通化、管理ツールの一元化、要員スキルの集約を実現。保守も単一窓口でワンストップ対応。",
         C["amber_500"]),
    ]

    card_h = Inches(1.0)
    half_w = CONTENT_WIDTH / 2 - Inches(0.1)
    for i, (icon, title, body, color) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = MARGIN_LEFT + col * (half_w + Inches(0.2))
        card_y = y + row * (card_h + Inches(0.1))
        _add_bullet_card(slide, x, card_y, half_w, card_h, title, body, color)


def slide_09_resource_storage(prs: Presentation) -> None:
    """スライド9: リソースPod向けストレージ構成。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 9, "リソースPod向けストレージ構成")
    y = _add_message_bar(slide, "VM領域に最適化したNFSストレージ構成で高性能・高可用性を実現")

    # 左: テキスト
    items = [
        ("提案機種", "NetApp AFF Aシリーズ / ONTAP 9.x\nプロトコル: NFS v4.1"),
        ("HA構成", "コントローラ冗長（HA Pair）\nディスク冗長（RAID-TEC）"),
        ("配置台数", "商用東京: 2台 / 大阪: 1台\nSTG: 2台 / 開発: 2台"),
        ("接続構成", "ESXi → Storage SW → NetApp\nNFS Datastore としてマウント"),
        ("サイジング", "現行SDN2実績値 × 成長率 × バッファ\n5年間の収容需要を十分にカバー"),
    ]

    for i, (title, body) in enumerate(items):
        _add_bullet_card(slide, MARGIN_LEFT, y + i * Inches(0.9),
                         Inches(5.5), Inches(0.8), title, body, C["blue_500"])


def slide_10_mgmt_storage(prs: Presentation) -> None:
    """スライド10: 管理Pod向けストレージ構成。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 10, "管理Pod向けストレージ構成")
    y = _add_message_bar(slide, "VM領域/データ領域/バックアップ領域を用途別に最適設計")

    headers = ["種別", "用途", "機種", "プロトコル", "容量目安", "冗長方式"]
    rows = [
        ["VM Storage", "管理VM用NFS", "NetApp AFF", "NFS", "XX TB", "HA Pair"],
        ["Data Storage", "ファイル共有/ログ保存", "NetApp AFF", "NFS/CIFS", "XX TB", "HA Pair"],
        ["Backup Storage", "SnapVault先", "NetApp FAS", "SnapVault", "XX TB", "HA Pair"],
    ]
    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.5),
                headers, rows, col_widths=[1.5, 2, 1.5, 1.2, 1, 1.2])

    # SnapMirror説明
    _add_bullet_card(slide, MARGIN_LEFT, y + Inches(2.5),
                     CONTENT_WIDTH, Inches(0.8),
                     "SnapMirror連携（SDN2⇔SDN3）",
                     "現行SDN2のONTAP 9.13.1との互換性を活かし、SnapMirrorによるミラーリングでデータ移行をシームレスに実現。移行完了後はDR用途に転用可能。",
                     C["teal_500"])


def slide_11_netapp_advantage(prs: Presentation) -> None:
    """スライド11: NetApp製品の優位性。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 11, "NetApp製品の優位性")
    y = _add_message_bar(slide, "NetApp ONTAPが通信事業者の仮想化基盤に最適な理由")

    items = [
        ("実績", "通信事業者の仮想化基盤への豊富な導入実績。SDN2でも採用されており、運用ノウハウの蓄積あり。", C["blue_500"]),
        ("統一管理", "ONTAPベースの統一アーキテクチャで運用効率を最大化。リソースPod/管理Pod/バックアップを共通の管理フレームワークで運用。", C["teal_500"]),
        ("データ保護", "SnapMirror/SnapVaultによる効率的なバックアップ・DR。ブロック単位の差分転送で帯域を節約しつつ確実にデータを保護。", C["green_500"]),
        ("移行容易性", "SDN2（現行ONTAP 9.13.1）からの親和性の高い移行パス。SnapMirrorで事前にデータ同期し、サービス影響を最小化。", C["amber_500"]),
        ("ライフサイクル", "5年以上の長期サポートとNDU（無停止アップグレード）。ONTAPのメジャーバージョンアップもサービス停止なしで実施可能。", C["blue_400"]),
    ]

    card_h = Inches(0.85)
    for i, (title, body, color) in enumerate(items):
        _add_bullet_card(slide, MARGIN_LEFT, y + i * (card_h + Inches(0.06)),
                         CONTENT_WIDTH, card_h, title, body, color)


def slide_12_resource_pod_design(prs: Presentation) -> None:
    """スライド12: リソースPod設計方針。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 12, "リソースPod設計方針")
    y = _add_message_bar(slide, "ESXiクラスタ＋NFSストレージの標準的な構成で信頼性と拡張性を両立")

    # 左: テキスト
    items = [
        ("vSphereクラスタ構成", "ESXi x18台でクラスタ構成。DRS/HA有効化。"),
        ("NFSデータストア", "ローカルストレージ未使用。全VMをNFS上に格納。"),
        ("AG物理障害分離", "ラック/電源系統をまたぐ障害分離設計。"),
        ("Pod独立設計", "Pod間でリソース混在なし。障害影響範囲を限定。"),
        ("VCF9対応", "SDDC Managerによるライフサイクル管理。"),
    ]
    for i, (title, body) in enumerate(items):
        _add_bullet_card(slide, MARGIN_LEFT, y + i * Inches(0.78),
                         Inches(5.0), Inches(0.7), title, body, C["blue_500"])

    # 右: 図
    _try_add_image(slide, Inches(6.5), y, Inches(5.8), Inches(4.0),
                   "slide12_resource_pod.drawio")


def slide_13_mgmt_pod_design(prs: Presentation) -> None:
    """スライド13: 管理Pod設計方針。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 13, "管理Pod設計方針")
    y = _add_message_bar(slide, "VCF9/vCSA/OCP/AAPを統合した管理基盤で運用自動化を実現")

    items = [
        ("VM基盤", "ESXiクラスタ上にvCSA/NSO/GitLab等の管理VMを集約。"),
        ("OCP基盤", "ベアメタル上にOpenShift。API-GW/AAP等をコンテナ化。"),
        ("ストレージ", "VM/Data/Backupの3種類を用途別に最適配置。"),
    ]
    for i, (title, body) in enumerate(items):
        _add_bullet_card(slide, MARGIN_LEFT, y + i * Inches(0.75),
                         Inches(5.0), Inches(0.65), title, body, C["green_500"])

    _try_add_image(slide, Inches(6.5), y, Inches(5.8), Inches(4.5),
                   "slide13_mgmt_pod.drawio")


def slide_14_mgmt_components(prs: Presentation) -> None:
    """スライド14: 管理VM・コンテナ設計方針。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 14, "管理VM・コンテナ設計方針")
    y = _add_message_bar(slide, "管理コンポーネントの適切な配置と冗長設計で運用基盤の安定性を確保")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.0),
                   "slide14_mgmt_components.drawio")


def slide_15_mgmt_nw(prs: Presentation) -> None:
    """スライド15: MgmtNW・運用サーバ設計方針。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 15, "MgmtNW・運用サーバ設計方針")
    y = _add_message_bar(slide, "管理NWと運用サーバの堅牢な設計でサービス可視性と障害対応力を確保")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.0),
                   "slide15_mgmt_nw.drawio")


def slide_16_migration(prs: Presentation) -> None:
    """スライド16: SDN2→SDN3移行設計。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 16, "SDN2→SDN3移行設計")
    y = _add_message_bar(slide, "移行自動化とミラーリング活用で顧客影響を最小化した段階的移行を実現")

    # 上: 移行フロー図
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(3.2),
                   "slide16_migration_flow.drawio")

    # 下: 主要ポイント
    bottom_y = y + Inches(3.4)
    items = [
        ("冗長あり VM", "vMotion: サービス影響 < 3分"),
        ("冗長なし VM", "Cold Migration: 影響 < 30分"),
        ("ファイルストレージ", "SnapMirror: 事前同期→カットオーバー"),
        ("自動化", "Ansible Playbook + カスタムスクリプト"),
    ]
    w = CONTENT_WIDTH / 4 - Inches(0.1)
    for i, (title, body) in enumerate(items):
        x = MARGIN_LEFT + i * (w + Inches(0.13))
        _add_bullet_card(slide, x, bottom_y, w, Inches(0.75), title, body, C["blue_500"])


def slide_17_performance(prs: Presentation) -> None:
    """スライド17: パフォーマンス。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 17, "パフォーマンス（性能）")
    y = _add_message_bar(slide, "サイジング根拠に基づき要求性能を十分に満たす構成を提案")

    headers = ["指標", "リソースPod", "管理Pod", "サイジング根拠"]
    rows = [
        ["IOPS", "XX,XXX IOPS", "XX,XXX IOPS", "現行SDN2実績値 × 1.3"],
        ["スループット", "XX GB/s", "XX GB/s", "NFS v4.1 マルチパス"],
        ["レイテンシ", "< X ms", "< X ms", "All Flash (SSD)"],
        ["NW帯域 (サービス)", "25GbE x2 (LAG)", "25GbE x2 (LAG)", "Leaf-Server間"],
        ["NW帯域 (ストレージ)", "25GbE x2 (LAG)", "25GbE x2 (LAG)", "StSW-Storage間"],
    ]
    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.45),
                headers, rows, col_widths=[2, 1.5, 1.5, 3])


def slide_18_capacity(prs: Presentation) -> None:
    """スライド18: スケール（収容量）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 18, "スケール（収容量）")
    y = _add_message_bar(slide, "現行需要の収容に加え、5年分の成長を見据えた十分な余裕を確保")

    headers = ["項目", "初期（構築時）", "5年後予測", "上限", "余裕率"]
    rows = [
        ["VM収容数（1Pod）", "XXX VM", "XXX VM", "XXX VM", "XX%"],
        ["VM収容数（全体）", "X,XXX VM", "X,XXX VM", "X,XXX VM", "XX%"],
        ["ストレージ容量（1Pod）", "XX TB", "XX TB", "XX TB", "XX%"],
        ["ストレージ容量（全体）", "XXX TB", "XXX TB", "XXX TB", "XX%"],
        ["顧客数", "XXX 顧客", "XXX 顧客", "XXX 顧客", "XX%"],
    ]
    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.45),
                headers, rows, col_widths=[2.5, 1.5, 1.5, 1.5, 1])


def slide_19_scalability(prs: Presentation) -> None:
    """スライド19: スケーラビリティ。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 19, "スケーラビリティ（拡張性）")
    y = _add_message_bar(slide, "Pod単位の拡張設計でサービス無停止の柔軟な増設を実現")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.0),
                   "slide19_scalability.drawio")


def slide_20_availability(prs: Presentation) -> None:
    """スライド20: 可用性設計。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 20, "可用性設計")
    y = _add_message_bar(slide, "全コンポーネントの冗長構成で単一障害点を排除し、1分以内の自動切替を実現")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.0),
                   "slide20_availability.drawio")


def slide_21_backup(prs: Presentation) -> None:
    """スライド21: バックアップ・リストア方式。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 21, "バックアップ・リストア方式")
    y = _add_message_bar(slide, "コンポーネントごとに最適なバックアップ方式を選定し、迅速なリストアを保証")

    headers = ["対象", "方式", "頻度", "保持期間", "RTO", "RPO"]
    rows = [
        ["VM", "Snapshot + NetApp Snapshot", "日次", "7世代", "< 1h", "24h"],
        ["コンフィグ", "GitLab/NetBox", "変更時", "全履歴", "< 30min", "0"],
        ["ファイルデータ", "SnapVault", "日次", "30世代", "< 4h", "24h"],
        ["OCP/コンテナ", "etcd backup + PV snapshot", "日次", "7世代", "< 2h", "24h"],
        ["DR (東京→大阪)", "SnapMirror", "非同期", "常時", "< 4h", "RPO設定依存"],
    ]
    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.45),
                headers, rows, col_widths=[1.5, 2.5, 1, 1, 1, 1.5])


def slide_22_security(prs: Presentation) -> None:
    """スライド22: セキュリティ対応。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 22, "セキュリティ対応")
    y = _add_message_bar(slide, "セキュリティガイドラインの必須要件を全て満たし、多層防御を実現")

    headers = ["カテゴリ", "対応方針", "ツール/手段"]
    rows = [
        ["ユーザ認証", "当社認証サーバ連携 + SSHキーローカル認証", "LDAP/RADIUS"],
        ["エンドポイント", "Deep Security導入（導入不可サーバは理由明記）", "Trend Micro DSM"],
        ["ログ転送", "SIEM連携（振る舞い/システム/操作ログ）", "syslog → 外部SIEM"],
        ["構成管理", "Rapid7/NetBox導入支援", "Rapid7 InsightVM"],
        ["ストレージ暗号化", "NetApp ONTAPの暗号化機能活用", "NVE/NAE"],
    ]
    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.45),
                headers, rows, col_widths=[1.5, 3.5, 2])


def slide_23_equipment(prs: Presentation) -> None:
    """スライド23: 機器・ライセンス一覧。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 23, "機器・ライセンス一覧")
    y = _add_message_bar(slide, "当社調達対象の全機器・ライセンスをスペック・数量とともに明示")

    headers = ["品目", "モデル/製品名", "商用東京", "商用大阪", "STG", "開発", "合計"]
    rows = [
        ["NFS Storage (リソース)", "NetApp AFF AXXX", "2", "1", "2", "2", "7"],
        ["VM Storage (管理)", "NetApp AFF AXXX", "1", "—", "1", "1", "3"],
        ["Data Storage (管理)", "NetApp AFF AXXX", "1", "—", "1", "1", "3"],
        ["Backup Storage", "NetApp FAS XXXX", "1", "—", "1", "1", "3"],
        ["NSO", "Cisco NSO x.x", "2", "—", "2", "2", "6"],
        ["ストレージ監視SW", "NetApp AIQUM", "1", "—", "1", "1", "3"],
        ["インテリジェントPDU", "XXX", "各ラック", "各ラック", "各ラック", "各ラック", "XX"],
        ["ケーブル類", "各種", "—", "—", "—", "—", "XX"],
    ]
    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.42),
                headers, rows, col_widths=[2.5, 2.5, 1, 1, 0.8, 0.8, 0.8])


def slide_24_rack(prs: Presentation) -> None:
    """スライド24: ラック構成・電力設計。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 24, "ラック構成・電力設計")
    y = _add_message_bar(slide, "ファシリティ制約を遵守したラック設計と電力計算で安全な運用を実現")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.0),
                   "slide24_rack.drawio")


def slide_25_monitoring(prs: Presentation) -> None:
    """スライド25: ストレージモニタリング。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 25, "ストレージモニタリング")
    y = _add_message_bar(slide, "可視化された管理ツールでストレージの正常性・性能をリアルタイムに把握")

    items = [
        ("管理ツール概要", "NetApp Active IQ Unified Manager (AIQUM) を管理Pod上のVMとして構築。\nWebベースのGUIでストレージの状態を統合管理。", C["blue_500"]),
        ("ハードウェア監視", "障害検知 → SNMP Trap/e-mail で即時通知。\n故障部位の視覚的特定が可能。", C["red_500"]),
        ("性能監視", "IOPS/Latency/Throughput のリアルタイム表示。\n閾値超過時のアラート通知機能。", C["teal_500"]),
        ("容量管理", "ボリューム/LUN 単位の使用量・予測トレンド表示。\n容量逼迫の事前検知とアラート。", C["green_500"]),
    ]

    card_h = Inches(0.95)
    for i, (title, body, color) in enumerate(items):
        _add_bullet_card(slide, MARGIN_LEFT, y + i * (card_h + Inches(0.08)),
                         CONTENT_WIDTH, card_h, title, body, color)


def slide_26_maintenance(prs: Presentation) -> None:
    """スライド26: 運用保守サービス概要。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 26, "運用保守サービス概要")
    y = _add_message_bar(slide, "統合保守サービスで単一窓口のワンストップ対応を実現")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.0),
                   "slide26_maintenance.drawio")


def slide_27_incident(prs: Presentation) -> None:
    """スライド27: 保守体制・障害対応フロー。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 27, "保守体制・障害対応フロー")
    y = _add_message_bar(slide, "24/365の障害受付体制と4時間以内駆けつけでサービスへの影響を最小化")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.0),
                   "slide27_incident_flow.drawio")


def slide_28_monthly(prs: Presentation) -> None:
    """スライド28: 月次運用・トレーニング。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 28, "月次運用・トレーニング")
    y = _add_message_bar(slide, "月次報告とハンズオントレーニングで当社運用チームの自立運用を支援")

    # 月次報告
    _add_bullet_card(slide, MARGIN_LEFT, y,
                     Inches(5.5), Inches(1.8),
                     "月次運用定例",
                     "■ 実施: 月1回（リモートまたは指定場所）\n■ 報告資料:\n  ・障害/問い合わせ状況一覧\n  ・HW/SW不具合・脆弱性情報\n  ・EoS/EoL情報\n■ セキュリティ情報: 重大不具合時は即時提供",
                     C["blue_500"])

    # トレーニング
    _add_bullet_card(slide, MARGIN_LEFT, y + Inches(2.0),
                     Inches(5.5), Inches(1.5),
                     "トレーニング計画",
                     "■ 座学: 製品概要/アーキテクチャ\n■ ハンズオン: 実機操作/一次切り分け\n■ マニュアル: 運用手順書/障害対応手順書を提供",
                     C["green_500"])

    # 保守期間
    _add_bullet_card(slide, Inches(6.5), y,
                     Inches(5.5), Inches(1.2),
                     "保守契約条件",
                     "■ 保守期間: 検収後5年間\n■ 言語: 日本語\n■ 受付: 24/365（障害）/ 平日9-18（問い合わせ）\n■ 保守移管: 受入対応可能",
                     C["teal_500"])


def slide_29_schedule(prs: Presentation) -> None:
    """スライド29: 実施スケジュール。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 29, "実施スケジュール")
    y = _add_message_bar(slide, "RFP記載の納期を遵守し、主要マイルストーンを明確化した実行計画")

    # ガントチャート
    months = ["2025\n12月", "2026\n1月", "2月", "3月", "4月", "5月", "6月",
              "7月", "8月", "9月", "10月", "11月", "12月", "2027\n1月", "2月", "3月", "4月"]
    phases = ["発注・KO", "設計", "機器搬入", "開発環境構築", "テスト",
              "STG構築", "本番構築", "受入試験", "リリース"]

    n_rows = len(phases) + 1
    n_cols = len(months) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, MARGIN_LEFT, y,
                                       CONTENT_WIDTH, Inches(4.2))
    tbl = tbl_shape.table

    # ヘッダー行
    cell0 = tbl.cell(0, 0)
    cell0.text = "フェーズ"
    cell0.fill.solid()
    cell0.fill.fore_color.rgb = C["navy_800"]
    for p in cell0.text_frame.paragraphs:
        for r in p.runs:
            _set_font(r, FONT_HEADING, Pt(8), True, C["white"])

    for ci, m in enumerate(months):
        cell = tbl.cell(0, ci + 1)
        cell.text = m
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["navy_800"]
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _set_font(r, FONT_HEADING, Pt(7), True, C["white"])

    # タスク行
    for ri, phase in enumerate(phases):
        cell = tbl.cell(ri + 1, 0)
        cell.text = phase
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                _set_font(r, FONT_BODY, Pt(8), True, C["text_primary"])

    # ガントバー色設定 (row, start_col, end_col, color)
    gantt_bars = [
        (1, 1, 2, C["text_muted"]),        # 発注KO: 12月-1月
        (2, 2, 5, C["blue_500"]),           # 設計: 1月-4月
        (3, 4, 4, C["navy_700"]),           # 機器搬入: 3月
        (4, 5, 6, C["teal_500"]),           # 開発環境: 4月-5月
        (5, 7, 13, C["amber_500"]),         # テスト: 6月-12月
        (6, 7, 10, C["green_500"]),         # STG: 6月-9月
        (7, 11, 14, C["red_500"]),          # 本番: 10月-1月
        (8, 15, 15, C["navy_700"]),         # 受入: 2月
        (9, 16, 17, C["text_muted"]),       # リリース: 3月-4月
    ]

    for row, start, end, color in gantt_bars:
        for col in range(start, end + 1):
            if col < n_cols:
                _make_gantt_cell_colored(tbl, row, col, color)


def slide_30_project_org(prs: Presentation) -> None:
    """スライド30: プロジェクト体制。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 30, "プロジェクト体制")
    y = _add_message_bar(slide, "経験豊富なメンバーによる堅実なプロジェクト体制でリスクを最小化")
    _try_add_image(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(5.0),
                   "slide30_project_org.drawio")


def slide_31_cost(prs: Presentation) -> None:
    """スライド31: 費用構成概要。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 31, "費用構成概要")
    y = _add_message_bar(slide, "費用の全体構成を大分類で提示（詳細は別紙7に記載）")

    headers = ["費用大分類", "内訳", "備考"]
    rows = [
        ["機器・ライセンス費用", "ストレージ / NSO / PDU / ケーブル / 監視SW", "ストレージ含む/含まない2パターン"],
        ["導入費用", "設計 / 構築 / 試験 / 移行 / PM", "環境別に工数算出"],
        ["HW保守費用（5年）", "ストレージ / NSO / PDU のHW保守", "検収後5年間"],
        ["SW保守費用（5年）", "ONTAP / NSO / 監視SW のSW保守", "検収後5年間"],
    ]
    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.5),
                headers, rows, col_widths=[2, 3.5, 3])

    note_y = y + Inches(3.0)
    _add_bullet_card(slide, MARGIN_LEFT, note_y,
                     CONTENT_WIDTH, Inches(0.7),
                     "費用に関する注記",
                     "■ 値引きはコンポーネント毎に計上（出精値引き不可）\n■ 詳細は別紙7_情報回答フォーマットをご参照ください",
                     C["amber_500"])


def slide_32_assumptions(prs: Presentation) -> None:
    """スライド32: 前提条件・制約条件。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 32, "前提条件・制約条件")
    y = _add_message_bar(slide, "提案の前提事項と制約事項を明確化し、認識齟齬を防止")

    # 前提条件
    _add_bullet_card(slide, MARGIN_LEFT, y,
                     Inches(5.5), Inches(2.5),
                     "前提条件",
                     "■ 当社調達機器は別紙7記載の構成に基づく\n■ VMware/Red Hat/ゲストOSライセンスはSB様にて手配\n■ インストールメディアはSB様より貸与\n■ 本番環境作業端末はSB様提供\n■ VCFバージョンはSB様と協議の上決定",
                     C["blue_500"])

    # 制約条件
    _add_bullet_card(slide, Inches(6.5), y,
                     Inches(5.5), Inches(2.5),
                     "制約条件",
                     "■ 大阪拠点はリソースPodのみ（管理Podなし）\n■ 監視サーバの構築は本スコープ外\n  （SB様調達・設定投入）\n■ 移行設計対象（※5）は本構築スコープ外\n\n■ SB様側で手配・準備が必要な事項\n  ・サーバ/NW機器の調達\n  ・VMware/Red Hatライセンスの調達\n  ・DC設備（電源/冷却/ケーブリング）の準備",
                     C["red_500"])


def slide_33_summary(prs: Presentation) -> None:
    """スライド33: 本提案のまとめ。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 33, "本提案のまとめ")
    y = _add_message_bar(slide, "信頼性・拡張性・運用効率の3本柱で貴社SDN3基盤の成功を支える")

    cards = [
        ("1", "信頼性",
         "NetApp ONTAPの実績ある高可用性ストレージ基盤。HA Pair + RAID-TEC + SnapMirrorで鉄壁のデータ保護を実現。",
         C["blue_500"]),
        ("2", "拡張性",
         "Pod単位の水平拡張で5年先を見据えたスケーラブル設計。サービス無停止での柔軟な増設が可能。",
         C["teal_500"]),
        ("3", "運用効率",
         "統合保守＋自動化で運用負荷を軽減。24/365対応の単一窓口と月次報告でSB様の運用チームを全面支援。",
         C["green_500"]),
    ]

    card_h = Inches(1.2)
    for i, (icon, title, body, color) in enumerate(cards):
        _add_bullet_card(slide, MARGIN_LEFT, y + i * (card_h + Inches(0.1)),
                         CONTENT_WIDTH, card_h, title, body, color, icon)

    # 結びの言葉
    closing_y = y + 3 * (card_h + Inches(0.1)) + Inches(0.1)
    _add_textbox(slide, MARGIN_LEFT, closing_y, CONTENT_WIDTH, Inches(0.4),
                 "当社の強み: 通信事業者向け仮想化基盤の豊富な構築実績とNetApp製品の専門知見",
                 font_name=FONT_HEADING, font_size=Pt(13),
                 bold=True, color=C["navy_800"],
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, MARGIN_LEFT, closing_y + Inches(0.45), CONTENT_WIDTH, Inches(0.35),
                 "ぜひご採用をご検討いただけますようお願い申し上げます。",
                 font_size=Pt(12), color=C["text_secondary"],
                 alignment=PP_ALIGN.CENTER)


def slide_34_appendix(prs: Presentation) -> None:
    """スライド34: 別紙一覧。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, 34, "別紙一覧")
    y = _add_message_bar(slide, "本提案書の付属資料一覧")

    headers = ["別紙番号", "資料名", "内容"]
    rows = [
        ["別紙2", "要求対応表_設計要件（ストレージ）", "ストレージに関する設計要件への回答"],
        ["別紙3", "要求対応表_設計要件（リソース/管理Pod）", "Pod構成に関する設計要件への回答"],
        ["別紙4", "要求対応表_設計要件（管理VM/コンテナ）", "管理ワークロードに関する設計要件への回答"],
        ["別紙5", "要求対応表_設計要件（MgmtNW/運用SV）", "管理NW/運用サーバに関する設計要件への回答"],
        ["別紙6", "要求対応表", "その他要件への回答"],
        ["別紙7", "情報回答フォーマット", "価格一覧・作業工数を含む"],
        ["別紙10", "セキュリティガイドライン対応", "セキュリティ要件への対応状況"],
        ["別紙11", "SCM・リスク管理に関する確認", "サプライチェーン/リスク管理の回答"],
        ["別紙12", "見積・購買要件書", "見積条件・購買要件"],
    ]
    _make_table(slide, MARGIN_LEFT, y, CONTENT_WIDTH, Inches(0.4),
                headers, rows, col_widths=[1.2, 3.5, 4])


# ============================================================
# メイン
# ============================================================

def main() -> None:
    """全34スライドのPPTXを生成する。"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("=== SDN3 提案書 PPTX生成 ===")
    print()

    generators = [
        ("01", "表紙", slide_01_cover),
        ("02", "目次", slide_02_toc),
        ("03", "エグゼクティブサマリー", slide_03_executive_summary),
        ("04", "貴社要求の理解", slide_04_requirements),
        ("05", "商用環境全体構成", slide_05_commercial_overview),
        ("06", "STG/開発環境構成", slide_06_stg_dev_overview),
        ("07", "開発対象範囲", slide_07_scope),
        ("08", "ストレージ提案概要", slide_08_storage_overview),
        ("09", "リソースPodストレージ", slide_09_resource_storage),
        ("10", "管理Podストレージ", slide_10_mgmt_storage),
        ("11", "NetApp優位性", slide_11_netapp_advantage),
        ("12", "リソースPod設計", slide_12_resource_pod_design),
        ("13", "管理Pod設計", slide_13_mgmt_pod_design),
        ("14", "管理VM/コンテナ設計", slide_14_mgmt_components),
        ("15", "MgmtNW/運用SV設計", slide_15_mgmt_nw),
        ("16", "移行設計", slide_16_migration),
        ("17", "パフォーマンス", slide_17_performance),
        ("18", "スケール", slide_18_capacity),
        ("19", "スケーラビリティ", slide_19_scalability),
        ("20", "可用性設計", slide_20_availability),
        ("21", "バックアップ", slide_21_backup),
        ("22", "セキュリティ", slide_22_security),
        ("23", "機器/ライセンス", slide_23_equipment),
        ("24", "ラック/電力", slide_24_rack),
        ("25", "モニタリング", slide_25_monitoring),
        ("26", "保守サービス", slide_26_maintenance),
        ("27", "障害対応", slide_27_incident),
        ("28", "月次運用", slide_28_monthly),
        ("29", "スケジュール", slide_29_schedule),
        ("30", "プロジェクト体制", slide_30_project_org),
        ("31", "費用構成", slide_31_cost),
        ("32", "前提/制約条件", slide_32_assumptions),
        ("33", "まとめ", slide_33_summary),
        ("34", "別紙一覧", slide_34_appendix),
    ]

    for num, label, func in generators:
        print(f"  [{num}] {label}")
        func(prs)

    output_path = os.path.join(OUTPUT_DIR, "proposal.pptx")
    prs.save(output_path)
    print()
    print(f"=== 完了: {len(generators)}枚のスライドを生成 ===")
    print(f"出力: {output_path}")


if __name__ == "__main__":
    main()
